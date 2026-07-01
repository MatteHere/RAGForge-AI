from pathlib import Path

import fitz
from docx import Document as DocxDocument
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


class DocumentParser:
    def parse_document(self, file_path: str) -> str:
        path = Path(file_path)
        extension = path.suffix.lower()

        if extension not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {extension}")

        if extension == ".pdf":
            return self._parse_pdf(path)

        if extension == ".docx":
            return self._parse_docx(path)

        if extension == ".pptx":
            return self._parse_pptx(path)

        if extension in {".txt", ".md"}:
            return self._parse_text_file(path)

        raise ValueError(f"No parser available for: {extension}")

    def _parse_pdf(self, path: Path) -> str:
        text_parts = []

        with fitz.open(path) as document:
            for page in document:
                text_parts.append(page.get_text())

        return "\n".join(text_parts).strip()

    def _parse_docx(self, path: Path) -> str:
        document = DocxDocument(path)
        paragraphs = [paragraph.text for paragraph in document.paragraphs]

        return "\n".join(paragraphs).strip()

    def _parse_pptx(self, path: Path) -> str:
        presentation = Presentation(path)
        slide_text = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

        return "\n".join(slide_text).strip()

    def _parse_text_file(self, path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="ignore").strip()